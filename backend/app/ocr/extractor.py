"""
EXTRACTOR DE IMÁGENES
----------------------
Corresponde al nodo "EXTRACTOR DE IMÁGENES -> SOLO imágenes" del diagrama.

Cuando el usuario sube un DOCUMENTO (PDF/DOCX/PPTX/XLSX) en lugar de una
imagen individual, este módulo recorre el archivo y extrae únicamente
las imágenes incrustadas, descartando texto nativo, tablas, etc.
Cada imagen extraída sigue el mismo camino que una "imagen individual"
(preprocesamiento -> OCR router -> OCR).

Se aplican dos filtros de calidad tras la extracción:
- Se descartan imágenes muy pequeñas (íconos, viñetas, logos) que no
  son contenido real a transcribir.
- Se deduplican imágenes idénticas (un mismo logo repetido en varias
  diapositivas/páginas no debería generar N resultados idénticos).
"""

import hashlib
import io
from typing import List

import fitz  # PyMuPDF
from PIL import Image
from docx import Document as DocxDocument
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import openpyxl

from ..config import settings


def _bytes_to_pil(data: bytes) -> Image.Image:
    return Image.open(io.BytesIO(data)).convert("RGB")


def extract_from_pdf(path: str) -> List[Image.Image]:
    """Extrae imágenes embebidas de cada página de un PDF."""
    imagenes: List[Image.Image] = []
    doc = fitz.open(path)
    try:
        for page_index in range(len(doc)):
            page = doc[page_index]
            for img in page.get_images(full=True):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    imagenes.append(_bytes_to_pil(base_image["image"]))
                except Exception:
                    continue

        # Si el PDF no tiene imágenes embebidas (p.ej. es un PDF escaneado
        # renderizado como página completa), renderizamos cada página como
        # imagen para que igual pueda pasar por el pipeline de OCR.
        if not imagenes:
            for page_index in range(len(doc)):
                page = doc[page_index]
                pix = page.get_pixmap(dpi=200)
                imagenes.append(_bytes_to_pil(pix.tobytes("png")))
    finally:
        doc.close()
    return imagenes


def extract_from_docx(path: str) -> List[Image.Image]:
    """Extrae imágenes embebidas de un documento Word."""
    imagenes: List[Image.Image] = []
    doc = DocxDocument(path)
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                imagenes.append(_bytes_to_pil(rel.target_part.blob))
            except Exception:
                continue
    return imagenes


def _extraer_fotos_de_shapes(shapes) -> List[Image.Image]:
    """Recorre shapes de una diapositiva de PowerPoint, incluyendo
    shapes dentro de grupos (recursivo), y extrae solo las imágenes."""
    imagenes: List[Image.Image] = []
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                imagenes.extend(_extraer_fotos_de_shapes(shape.shapes))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                imagenes.append(_bytes_to_pil(shape.image.blob))
        except Exception:
            continue
    return imagenes


def extract_from_pptx(path: str) -> List[Image.Image]:
    """Extrae imágenes embebidas de una presentación PowerPoint,
    incluyendo imágenes dentro de grupos de objetos."""
    imagenes: List[Image.Image] = []
    prs = Presentation(path)
    for slide in prs.slides:
        imagenes.extend(_extraer_fotos_de_shapes(slide.shapes))
    return imagenes


def extract_from_xlsx(path: str) -> List[Image.Image]:
    """Extrae imágenes embebidas de un libro Excel."""
    imagenes: List[Image.Image] = []
    wb = openpyxl.load_workbook(path)
    for ws in wb.worksheets:
        for image in getattr(ws, "_images", []):
            try:
                blob = image._data()
                imagenes.append(_bytes_to_pil(blob))
            except Exception:
                continue
    return imagenes


EXTRACTORS = {
    ".pdf": extract_from_pdf,
    ".docx": extract_from_docx,
    ".pptx": extract_from_pptx,
    ".xlsx": extract_from_xlsx,
}


def _filtrar_y_deduplicar(imagenes: List[Image.Image]) -> List[Image.Image]:
    """Descarta imágenes demasiado pequeñas (íconos/viñetas/logos) y
    elimina duplicados exactos (misma imagen repetida en el documento)."""
    vistas = set()
    resultado: List[Image.Image] = []
    min_size = settings.MIN_EXTRACTED_IMAGE_SIZE

    for imagen in imagenes:
        ancho, alto = imagen.size
        if ancho < min_size or alto < min_size:
            continue

        huella = hashlib.md5(imagen.tobytes()).hexdigest()
        if huella in vistas:
            continue
        vistas.add(huella)
        resultado.append(imagen)

    return resultado


def extract_images(path: str, extension: str) -> List[Image.Image]:
    """Punto de entrada único: recibe la ruta del documento y su extensión,
    devuelve la lista de imágenes PIL de contenido real (sin íconos
    diminutos ni duplicados) contenidas en él."""
    extractor = EXTRACTORS.get(extension.lower())
    if not extractor:
        raise ValueError(f"Extensión de documento no soportada: {extension}")

    imagenes_crudas = extractor(path)
    return _filtrar_y_deduplicar(imagenes_crudas)
